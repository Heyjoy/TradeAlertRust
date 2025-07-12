#!/usr/bin/env python3
"""
Office 文档读取工具
用于读取 .docx 和 .pptx 文件内容
"""

import sys
import os
from pathlib import Path
from docx import Document
from pptx import Presentation

def read_docx(file_path):
    """读取 Word 文档内容"""
    try:
        doc = Document(file_path)
        content = []
        
        print(f"\n=== {Path(file_path).name} ===")
        print(f"段落数量: {len(doc.paragraphs)}")
        print(f"表格数量: {len(doc.tables)}")
        
        # 读取段落
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                content.append(f"段落 {i+1}: {para.text}")
        
        # 读取表格
        for i, table in enumerate(doc.tables):
            content.append(f"\n表格 {i+1}:")
            for row_idx, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                if any(row_data):  # 只显示非空行
                    content.append(f"  行 {row_idx+1}: {' | '.join(row_data)}")
        
        return "\n".join(content)
        
    except Exception as e:
        return f"读取 Word 文档失败: {e}"

def read_pptx(file_path):
    """读取 PowerPoint 文档内容"""
    try:
        prs = Presentation(file_path)
        content = []
        
        print(f"\n=== {Path(file_path).name} ===")
        print(f"幻灯片数量: {len(prs.slides)}")
        
        for i, slide in enumerate(prs.slides):
            content.append(f"\n幻灯片 {i+1}:")
            
            # 读取文本框内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(f"  - {shape.text}")
                
                # 读取表格内容
                if shape.has_table:
                    content.append(f"  表格:")
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        if any(row_data):
                            content.append(f"    行 {row_idx+1}: {' | '.join(row_data)}")
        
        return "\n".join(content)
        
    except Exception as e:
        return f"读取 PowerPoint 文档失败: {e}"

def main():
    """主函数"""
    if len(sys.argv) < 2:
        print("用法: python3 read_office_docs.py <file_path_or_directory>")
        sys.exit(1)
    
    path = Path(sys.argv[1])
    
    if path.is_file():
        # 处理单个文件
        if path.suffix.lower() == '.docx':
            content = read_docx(path)
        elif path.suffix.lower() == '.pptx':
            content = read_pptx(path)
        else:
            print(f"不支持的文件类型: {path.suffix}")
            sys.exit(1)
        
        print(content)
        
    elif path.is_dir():
        # 处理目录中的所有 Office 文档
        office_files = []
        office_files.extend(path.glob("**/*.docx"))
        office_files.extend(path.glob("**/*.pptx"))
        
        if not office_files:
            print(f"在目录 {path} 中未找到 Office 文档")
            sys.exit(1)
        
        all_content = []
        
        for file_path in sorted(office_files):
            print(f"\n处理文件: {file_path}")
            
            if file_path.suffix.lower() == '.docx':
                content = read_docx(file_path)
            elif file_path.suffix.lower() == '.pptx':
                content = read_pptx(file_path)
            else:
                continue
            
            all_content.append(content)
            print(content)
        
        # 保存到文件
        output_file = path / "extracted_content.txt"
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write("\n\n" + "="*50 + "\n\n".join(all_content))
        
        print(f"\n所有内容已保存到: {output_file}")
    
    else:
        print(f"路径不存在: {path}")
        sys.exit(1)

if __name__ == "__main__":
    main()